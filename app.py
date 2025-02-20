import os
import torch
import logging
import gradio as gr
import whisper
from transformers import pipeline
from pydub import AudioSegment
from pptx import Presentation
from PyPDF2 import PdfReader

# Set up logging
logging.basicConfig(level=logging.INFO)

# Summarization pipeline using a pre-trained Transformer
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", tokenizer="facebook/bart-large-cnn")

def preprocess_audio_file(audio_file_path: str) -> str:
    """
    Converts the input audio file to WAV format with 16kHz sample rate and mono channel.
    """
    output_wav_file = f"{os.path.splitext(audio_file_path)[0]}_converted.wav"
    try:
        audio = AudioSegment.from_file(audio_file_path)
        audio = audio.set_frame_rate(16000).set_channels(1)
        audio.export(output_wav_file, format="wav")
    except Exception as e:
        logging.error(f"Error during audio preprocessing: {e}")
        raise RuntimeError("Failed to preprocess audio.")
    return output_wav_file

def extract_text_from_ppt(ppt_file: str) -> str:
    """
    Extracts text from PowerPoint slides.
    """
    text = ""
    try:
        presentation = Presentation(ppt_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from PPT: {e}")
        raise RuntimeError("Failed to process PPT.")
    return text.strip()

def extract_text_from_pdf(pdf_file: str) -> str:
    """
    Extracts text from PDF slides.
    """
    text = ""
    try:
        reader = PdfReader(pdf_file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from PDF: {e}")
        raise RuntimeError("Failed to process PDF.")
    return text.strip()

def summarize_with_dynamic_window(text: str, max_summary_length: int = 200) -> str:
    """
    Summarizes long transcripts using a dynamic sliding window for chunking.
    """
    try:
        if not text or len(text.strip()) < 50:
            return "Transcript is too short to summarize."

        chunk_size = 1500
        overlap = 200
        summaries = []
        start = 0

        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunk = text[start:end]
            summary = summarizer(chunk, max_length=max_summary_length, min_length=50, do_sample=False)
            summaries.append(summary[0]['summary_text'])
            start += chunk_size - overlap
            chunk_size = min(chunk_size + 500, 3000)

        final_summary = " ".join(summaries)
        return final_summary

    except Exception as e:
        logging.error(f"Error during summarization: {e}")
        return "Error summarizing text."

def translate_summarize_and_include_slides(audio_file: str, context: str, whisper_model_name: str, slide_file: str = None) -> tuple[str, str]:
    """
    Transcribes the audio, optionally extracts slide text, and generates a combined summary.
    """
    if not audio_file:
        return "No audio file provided.", None

    try:
        model = whisper.load_model(whisper_model_name, device="cpu")
    except Exception as e:
        logging.error(f"Error loading Whisper model: {e}")
        return f"Error loading Whisper model: {e}", None

    try:
        audio_file_wav = preprocess_audio_file(audio_file)
        result = model.transcribe(audio_file_wav)
        transcript = result.get("text", "")
        if not transcript:
            return "Transcription failed or no text found in the audio.", None
    except Exception as e:
        logging.error(f"Error during transcription: {e}")
        return f"Error during transcription: {e}", None
    finally:
        if os.path.exists(audio_file_wav):
            os.remove(audio_file_wav)

    # Extract slide text if provided
    slide_text = ""
    if slide_file:
        if slide_file.endswith(".pptx"):
            slide_text = extract_text_from_ppt(slide_file) or ""
        elif slide_file.endswith(".pdf"):
            slide_text = extract_text_from_pdf(slide_file) or ""
        else:
            return "Unsupported slide file type.", None

    # Combine transcript and slide text
    combined_text = f"{slide_text}\n\n{transcript}" if slide_text else transcript
    if not combined_text.strip():
        return "No text available for summarization.", None

    # Summarize the combined text
    summary = summarize_with_dynamic_window(combined_text)
    transcript_file = "transcript.txt"
    with open(transcript_file, "w") as f:
        f.write(transcript)

    return summary, transcript_file

# Gradio interface
def gradio_app(audio, context: str, whisper_model_name: str, slide_file: str = None) -> tuple[str, str]:
    """
    Gradio app to summarize meetings with optional slide input.
    """
    return translate_summarize_and_include_slides(audio, context, whisper_model_name, slide_file)

if __name__ == "__main__":
    iface = gr.Interface(
        fn=gradio_app,
        inputs=[
            gr.Audio(type="filepath", label="Upload an audio file"),  # Use "filepath" or "numpy" for type
            gr.Textbox(label="Context (optional)", placeholder="Provide any additional context"),
            gr.Textbox(label="Enter Whisper model name", placeholder="e.g., base, large-v3"),
            gr.File(label="Upload Slides (optional, PPT/PDF)", type="filepath")
        ],
        outputs=[
            gr.Textbox(label="Summary", show_copy_button=True),
            gr.File(label="Download Transcript")
        ],
        title="Meeting Summarizer with Slide Integration",
        description="Upload an audio file and optionally upload slides (PPT or PDF) to get a summarized transcript.",
    )

    iface.launch(debug=True, server_port=7860, server_name="0.0.0.0", share=True)
